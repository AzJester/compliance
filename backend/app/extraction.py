from __future__ import annotations

import io
import re
from dataclasses import dataclass
from pathlib import Path

import fitz
from docx import Document as WordDocument
from openpyxl import load_workbook
from pptx import Presentation

from .models import DocumentStatus


@dataclass(frozen=True, slots=True)
class ExtractionResult:
    text: str
    status: DocumentStatus
    error: str | None = None

    @property
    def character_count(self) -> int:
        return len(self.text)


def _clean_text(parts: list[str]) -> str:
    return "\n".join(part.strip() for part in parts if part and part.strip()).strip()


def _extract_pdf(data: bytes) -> ExtractionResult:
    with fitz.open(stream=data, filetype="pdf") as document:
        parts: list[str] = []
        for index, page in enumerate(document, start=1):
            page_text = page.get_text("text").strip()
            if page_text:
                parts.extend((f"[PDF Page {index}]", page_text))
        text = _clean_text(parts)
    if not text:
        return ExtractionResult(text="", status=DocumentStatus.NEEDS_OCR)
    return ExtractionResult(text=text, status=DocumentStatus.EXTRACTED)


def _extract_docx(data: bytes) -> ExtractionResult:
    document = WordDocument(io.BytesIO(data))
    parts: list[str] = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        if paragraph.text.strip():
            parts.extend((f"[DOCX Paragraph {index}]", paragraph.text))
    for table_index, table in enumerate(document.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            row_text = "\t".join(cell.text for cell in row.cells)
            if row_text.strip():
                parts.extend((f"[DOCX Table {table_index} Row {row_index}]", row_text))
    return ExtractionResult(text=_clean_text(parts), status=DocumentStatus.EXTRACTED)


def _cell_text(cell: object) -> str:
    value = cell.value  # type: ignore[attr-defined]
    if value is None:
        return ""
    number_format = str(getattr(cell, "number_format", ""))
    if (
        isinstance(value, (int, float))
        and not isinstance(value, bool)
        and float(value).is_integer()
        and re.fullmatch(r"0+", number_format)
    ):
        return f"{int(value):0{len(number_format)}d}"
    return str(value)


def _extract_xlsx(data: bytes) -> ExtractionResult:
    # keep_links=False prevents openpyxl from retaining external workbook packages;
    # no supported extractor performs network I/O.
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True, keep_links=False)
    parts: list[str] = []
    try:
        for worksheet in workbook.worksheets:
            for row_index, row in enumerate(worksheet.iter_rows(), start=1):
                values = [_cell_text(cell) for cell in row]
                if any(values):
                    parts.extend(
                        (
                            f'[XLSX Sheet "{worksheet.title}" Row {row_index}]',
                            "\t".join(values),
                        )
                    )
    finally:
        workbook.close()
    return ExtractionResult(text=_clean_text(parts), status=DocumentStatus.EXTRACTED)


def _extract_pptx(data: bytes) -> ExtractionResult:
    presentation = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[Slide {index}]")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append("\t".join(cell.text for cell in row.cells))
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                parts.append(notes_frame.text)
    return ExtractionResult(text=_clean_text(parts), status=DocumentStatus.EXTRACTED)


def extract_document(relative_path: str, data: bytes) -> ExtractionResult:
    """Extract local text and reduce parser failures to content-safe user messages."""

    extension = Path(relative_path.split("!/")[-1]).suffix.lower()
    if extension == ".zip":
        return ExtractionResult(text="", status=DocumentStatus.ARCHIVE_EXPANDED)

    extractors = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".pptx": _extract_pptx,
    }
    try:
        result = extractors[extension](data)
        if result.status == DocumentStatus.EXTRACTED and not result.text.strip():
            return ExtractionResult(
                text="",
                status=DocumentStatus.ERROR,
                error=(
                    "The document was preserved, but no searchable text could be extracted. "
                    "Export it as a searchable PDF or text-based Office document, then upload "
                    "it again."
                ),
            )
        return result
    except Exception:
        return ExtractionResult(
            text="",
            status=DocumentStatus.ERROR,
            error="The document was preserved, but its text could not be extracted.",
        )
