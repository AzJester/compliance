from __future__ import annotations

import io
import struct
import zipfile
from collections.abc import Callable, Iterator
from pathlib import Path

import fitz
import pytest
from docx import Document as WordDocument
from fastapi import FastAPI
from fastapi.testclient import TestClient
from openpyxl import Workbook
from pptx import Presentation

from backend.app.config import Settings
from backend.app.main import create_app


def make_test_client(app: FastAPI) -> TestClient:
    return TestClient(
        app,
        client=("127.0.0.1", 50_000),
        headers={"host": "127.0.0.1"},
    )


def pdf_bytes(text: str = "The offeror shall provide a management plan.") -> bytes:
    document = fitz.open()
    page = document.new_page()
    if text:
        page.insert_text((72, 72), text)
    data = document.tobytes()
    document.close()
    return data


def docx_bytes(text: str = "Section L requires a technical volume.") -> bytes:
    document = WordDocument()
    document.add_paragraph(text)
    stream = io.BytesIO()
    document.save(stream)
    return stream.getvalue()


def xlsx_bytes(text: str = "CDRL A001") -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Requirements"
    sheet.append(["Identifier", "Requirement"])
    sheet.append([text, "Submit monthly"])
    stream = io.BytesIO()
    workbook.save(stream)
    workbook.close()
    return stream.getvalue()


def pptx_bytes(text: str = "Evaluation factor: Technical Approach") -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Section M"
    slide.placeholders[1].text = text
    stream = io.BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def zip_bytes(entries: dict[str, bytes], *, compression: int = zipfile.ZIP_DEFLATED) -> bytes:
    stream = io.BytesIO()
    with zipfile.ZipFile(stream, "w", compression=compression) as archive:
        for name, data in entries.items():
            archive.writestr(name, data)
    return stream.getvalue()


def mark_zip_encrypted(data: bytes) -> bytes:
    """Set the encrypted bit in local and central headers for a test-only archive."""

    mutable = bytearray(data)
    offset = 0
    while offset < len(mutable) - 4:
        signature = bytes(mutable[offset : offset + 4])
        if signature == b"PK\x03\x04":
            flags = struct.unpack_from("<H", mutable, offset + 6)[0]
            struct.pack_into("<H", mutable, offset + 6, flags | 0x1)
        elif signature == b"PK\x01\x02":
            flags = struct.unpack_from("<H", mutable, offset + 8)[0]
            struct.pack_into("<H", mutable, offset + 8, flags | 0x1)
        offset += 1
    return bytes(mutable)


def mark_zip_unsupported_compression(data: bytes) -> bytes:
    """Replace supported ZIP methods with a synthetic unsupported method."""

    mutable = bytearray(data)
    offset = 0
    while offset < len(mutable) - 4:
        signature = bytes(mutable[offset : offset + 4])
        if signature == b"PK\x03\x04":
            struct.pack_into("<H", mutable, offset + 8, 99)
        elif signature == b"PK\x01\x02":
            struct.pack_into("<H", mutable, offset + 10, 99)
        offset += 1
    return bytes(mutable)


@pytest.fixture
def app_factory(tmp_path: Path) -> Callable[..., FastAPI]:
    counter = 0

    def factory(**overrides: int | str | Path | tuple[str, ...]) -> FastAPI:
        nonlocal counter
        counter += 1
        defaults: dict[str, int | str | Path | tuple[str, ...]] = {
            "data_dir": tmp_path / f"instance-{counter}",
            "host": "127.0.0.1",
            "port": 8000,
            "allowed_origins": (
                "http://127.0.0.1:8000",
                "http://localhost:8000",
                "http://127.0.0.1:5173",
                "http://localhost:5173",
            ),
        }
        defaults.update(overrides)
        return create_app(Settings(**defaults))  # type: ignore[arg-type]

    return factory


@pytest.fixture
def client(app_factory: Callable[..., FastAPI]) -> Iterator[TestClient]:
    with make_test_client(app_factory()) as local_client:
        yield local_client


@pytest.fixture
def project(client: TestClient) -> dict[str, object]:
    response = client.post(
        "/api/projects",
        json={
            "name": "Air Vehicle Support RFP",
            "solicitation_number": "FA0000-26-R-0001",
            "agency": "Department of the Air Force",
            "due_at": "2026-09-01T16:00:00-07:00",
            "due_timezone": "America/Phoenix",
            "sensitivity": "CUI",
        },
    )
    assert response.status_code == 201
    return response.json()
